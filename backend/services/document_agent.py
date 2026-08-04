"""services/document_agent.py — Document Agent (Document AI).

Upload PDF/DOCX/PPTX/XLSX/TXT → text is extracted and stored in SQLite
(the original file is NOT kept — privacy + zero clutter) → ask questions,
summarize, compare via Groq (free tier).

Voice path: `handle_voice_request()` powers the "ask my documents…",
"summarize the pdf about X", "search my documents for Y" commands.
"""

import os
import re
import sqlite3
import threading
import time
import uuid
from pathlib import Path

DB_PATH = Path(__file__).resolve().parent.parent / "data" / "documents.db"
_lock = threading.Lock()

MAX_UPLOAD_BYTES = int(os.getenv("FRIDAY_DOCUMENT_MAX_UPLOAD", str(10 * 1024 * 1024)))
MAX_TEXT_CHARS = 200_000

SUPPORTED_EXT = {
    ".pdf": "PDF document",
    ".docx": "Word document",
    ".pptx": "PowerPoint",
    ".xlsx": "Excel",
    ".txt": "Text",
    ".md": "Markdown",
}

_ASK_SYSTEM_PROMPT = (
    "You are F.R.I.D.A.Y.'s document intelligence engine. Answer the user's "
    "question using ONLY the provided document text. If the answer is not in "
    "the document, say so plainly. Keep answers concise (2-4 sentences)."
)

_SUMMARIZE_SYSTEM_PROMPT = (
    "You are F.R.I.D.A.Y.'s document intelligence engine. Summarize the "
    "document in 4-6 tight bullet points, starting each with '- '."
)

_COMPARE_SYSTEM_PROMPT = (
    "You are F.R.I.D.A.Y.'s document intelligence engine. Compare the two "
    "documents below: list 2-4 key similarities and 2-4 key differences, "
    "each as a bullet starting with '- '."
)


class DocumentUnavailableError(RuntimeError):
    """Raised when the LLM is unavailable or the document can't be processed."""


# ── DB ───────────────────────────────────────────────────────────────────

def init_documents_db():
    with _lock, _connect() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id         TEXT PRIMARY KEY,
                title      TEXT NOT NULL,
                ext        TEXT NOT NULL,
                text       TEXT NOT NULL,
                size       INTEGER DEFAULT 0,
                pages      INTEGER DEFAULT 0,
                created_at REAL NOT NULL
            )
        """)
        conn.execute("CREATE INDEX IF NOT EXISTS idx_documents_title ON documents(title)")
        conn.commit()


def _connect() -> sqlite3.Connection:
    DB_PATH.parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(DB_PATH, check_same_thread=False, timeout=10.0)
    conn.row_factory = sqlite3.Row
    return conn


# ── Text extraction ──────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> tuple[str, int]:
    from pypdf import PdfReader
    import io
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(parts), len(reader.pages)


def _extract_docx(data: bytes) -> tuple[str, int]:
    import docx
    import io
    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts), len(doc.paragraphs)


def _extract_pptx(data: bytes) -> tuple[str, int]:
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        parts.append("\n".join(slide_parts))
    return "\n".join(parts), len(prs.slides)


def _extract_xlsx(data: bytes) -> tuple[str, int]:
    import openpyxl
    import io
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            parts.append(f"--- Sheet: {ws.title} ---")
            parts.extend(rows)
    return "\n".join(parts), len(wb.sheetnames)


def _extract_txt(data: bytes) -> tuple[str, int]:
    for enc in ("utf-8", "latin-1"):
        try:
            return data.decode(enc), 1
        except UnicodeDecodeError:
            continue
    return "", 1


def extract_text(filename: str, data: bytes) -> tuple[str, dict]:
    """Return (text, meta). Raises DocumentUnavailableError for bad files."""
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXT:
        raise DocumentUnavailableError(
            f"Unsupported file type '{ext or '(none)'}'. Supported: "
            + ", ".join(sorted(SUPPORTED_EXT))
        )
    if ext == ".pdf":
        text, pages = _extract_pdf(data)
    elif ext == ".docx":
        text, pages = _extract_docx(data)
    elif ext == ".pptx":
        text, pages = _extract_pptx(data)
    elif ext == ".xlsx":
        text, pages = _extract_xlsx(data)
    else:
        text, pages = _extract_txt(data)

    text = text.strip()[:MAX_TEXT_CHARS]
    if not text:
        raise DocumentUnavailableError("No readable text found in that file (scanned PDFs aren't supported).")
    return text, {"ext": ext, "pages": pages, "size": len(data)}


# ── CRUD ─────────────────────────────────────────────────────────────────

def add_document(filename: str, data: bytes) -> dict:
    text, meta = extract_text(filename, data)
    title = Path(filename).stem.strip()[:150] or "Untitled document"
    doc = {
        "id": uuid.uuid4().hex,
        "title": title,
        "ext": meta["ext"],
        "text": text,
        "size": meta["size"],
        "pages": meta["pages"],
        "created_at": time.time(),
    }
    with _lock, _connect() as conn:
        conn.execute(
            "INSERT INTO documents (id, title, ext, text, size, pages, created_at) "
            "VALUES (?, ?, ?, ?, ?, ?, ?)",
            (doc["id"], doc["title"], doc["ext"], doc["text"], doc["size"], doc["pages"], doc["created_at"]),
        )
        conn.commit()
    return _public(doc)


def _public(doc: dict, include_text: bool = False) -> dict:
    out = {
        "id": doc["id"],
        "title": doc["title"],
        "ext": doc["ext"],
        "size": doc["size"],
        "pages": doc["pages"],
        "created_at": doc["created_at"],
        "snippet": doc["text"][:180].replace("\n", " ") if doc.get("text") else "",
    }
    if include_text:
        out["text"] = doc["text"]
    return out


def _row_public(row, include_text: bool = False) -> dict:
    return _public(dict(row), include_text=include_text)


def list_documents(limit: int = 50) -> list:
    with _lock, _connect() as conn:
        rows = conn.execute(
            "SELECT * FROM documents ORDER BY created_at DESC LIMIT ?", (limit,)
        ).fetchall()
    return [_row_public(r) for r in rows]


def get_document(doc_id: str, include_text: bool = True) -> dict | None:
    with _lock, _connect() as conn:
        row = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    return _row_public(row, include_text=include_text) if row else None


def delete_document(doc_id: str) -> bool:
    with _lock, _connect() as conn:
        cur = conn.execute("DELETE FROM documents WHERE id = ?", (doc_id,))
        conn.commit()
    return cur.rowcount > 0


def search_documents(query: str, limit: int = 10) -> list:
    q = (query or "").strip()
    if not q:
        return []
    with _lock, _connect() as conn:
        like = f"%{q}%"
        rows = conn.execute(
            """SELECT * FROM documents
               WHERE title LIKE ? OR text LIKE ?
               ORDER BY created_at DESC LIMIT ?""",
            (like, like, limit),
        ).fetchall()
    return [_row_public(r) for r in rows]


def find_document_by_keyword(query: str) -> dict | None:
    """Best title/content match for a voice request (or None)."""
    hits = search_documents(query, limit=5)
    if not hits:
        return None
    q = query.lower()
    # Prefer title matches, then earliest (most relevant content match)
    for h in hits:
        if q in h["title"].lower():
            return h
    return hits[0]


# ── LLM (Groq, free tier) ───────────────────────────────────────────────

def _llm(system_prompt: str, user_prompt: str, model: str | None = None) -> str:
    from services.brain import _get_groq_client

    client = _get_groq_client()
    if client is None:
        raise DocumentUnavailableError("GROQ_API_KEY is not configured — can't run document AI.")
    try:
        completion = client.chat.completions.create(
            model=model or os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile"),
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt[:30000]},
            ],
            temperature=0.2,
        )
        return (getattr(completion.choices[0].message, "content", "") or "").strip()
    except Exception as exc:
        raise DocumentUnavailableError(f"LLM call failed: {exc}") from exc


def ask_document(doc_id: str, question: str) -> str:
    doc = get_document(doc_id, include_text=True)
    if not doc:
        raise DocumentUnavailableError("Document not found.")
    question = (question or "").strip()
    if not question:
        raise DocumentUnavailableError("A question is required.")
    return _llm(
        _ASK_SYSTEM_PROMPT,
        f"DOCUMENT TITLE: {doc['title']}\n\nDOCUMENT TEXT:\n{doc['text'][:14000]}\n\n"
        f"QUESTION: {question}",
    )


def summarize_document(doc_id: str) -> str:
    doc = get_document(doc_id, include_text=True)
    if not doc:
        raise DocumentUnavailableError("Document not found.")
    return _llm(
        _SUMMARIZE_SYSTEM_PROMPT,
        f"DOCUMENT TITLE: {doc['title']}\n\nDOCUMENT TEXT:\n{doc['text'][:14000]}",
    )


def compare_documents(doc_ids: list) -> str:
    docs = [get_document(did, include_text=True) for did in doc_ids]
    docs = [d for d in docs if d]
    if len(docs) < 2:
        raise DocumentUnavailableError("Need at least two documents to compare.")
    parts = []
    for d in docs[:2]:
        parts.append(f"DOCUMENT '{d['title']}':\n{d['text'][:8000]}")
    return _llm(_COMPARE_SYSTEM_PROMPT, "\n\n----\n\n".join(parts))


# ── Voice path (used by brain.py) ────────────────────────────────────────

_VOICE_STRIPS = re.compile(
    r"^(?:hey\s+friday\s*)?(?:please\s+|can\s+you\s+|could\s+you\s+|will\s+you\s+)?"
    r"(?:ask|summarize|summarise|search|find|read)\s+(?:my\s+|the\s+|your\s+)?"
    r"(?:documents?|pdfs?|files?|docs?|resume|resumes?)\s+(?:about|for|on|regarding|in)\s*",
    re.IGNORECASE,
)


def handle_voice_request(text: str) -> str:
    """Handle 'ask my documents about X' / 'summarize the pdf about Y' etc."""
    lower = (text or "").lower()

    # Summarize request
    if re.search(r"\b(summariz|summaris)\w*\b", lower):
        q = re.sub(r"^(?:please\s+|can\s+you\s+)?(?:summariz\w*\s+)(?:my\s+|the\s+|your\s+)?(?:documents?|pdfs?|files?|docs?|resume|resumes?)\s*", "", lower).strip()
        doc = find_document_by_keyword(q or " ")
        if not doc:
            return "I couldn't find a matching document, Boss. Ask me to list your documents."
        return f"Summary of '{doc['title']}':\n" + summarize_document(doc["id"])

    # Search / list request
    if re.search(r"\b(search|find|list|show)\b", lower):
        q = _VOICE_STRIPS.sub("", lower).strip()
        results = search_documents(q) if q else list_documents(limit=8)
        if not results:
            return "No documents stored yet, Boss — upload one and I'll read it for you."
        lines = [f"Found {len(results)} document(s):"]
        lines += [f"- {d['title']} ({d['ext']}, {d['pages']} page(s))" for d in results[:8]]
        return " ".join(lines)

    # Default: ask
    q = _VOICE_STRIPS.sub("", lower).strip()
    doc = find_document_by_keyword(q) if q else None
    if not doc:
        return "I couldn't find a matching document. Ask me to list your documents, Boss."
    return f"From '{doc['title']}': " + ask_document(doc["id"], q)


# Ensure the table exists on first import.
init_documents_db()
